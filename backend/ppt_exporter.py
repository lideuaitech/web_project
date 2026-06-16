from pptx import Presentation
from pptx.util import Inches
from datetime import datetime
import pandas as pd
import uuid
import os

class PPTExporter:

    def __init__(self, df, summary):

        self.df = df
        self.summary = summary

    def export(self):

        prs = Presentation()

        # =========================
        # SLIDE 1 → TITLE
        # =========================

        slide_layout = prs.slide_layouts[0]

        slide = prs.slides.add_slide(
            slide_layout
        )

        title = slide.shapes.title

        subtitle = slide.placeholders[1]

        title.text = "Lideu AI Analytics Report"

        subtitle.text = (
            f"Generated on "
            f"{datetime.now().strftime('%d %b %Y %H:%M')}"
        )

        # =========================
        # SLIDE 2 → SUMMARY
        # =========================

        slide_layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(
            slide_layout
        )

        title = slide.shapes.title

        body = slide.placeholders[1]

        title.text = "AI Insights"

        body.text = self.summary

        # =========================
        # SLIDE 3 → TABLE
        # =========================

        slide_layout = prs.slide_layouts[5]

        slide = prs.slides.add_slide(
            slide_layout
        )

        title = slide.shapes.title

        title.text = "Analytics Data"

        rows = len(self.df) + 1

        cols = len(self.df.columns)

        table = slide.shapes.add_table(
            rows,
            cols,
            Inches(0.5),
            Inches(1.5),
            Inches(9),
            Inches(4)
        ).table

        # HEADERS
        for col_idx, column_name in enumerate(self.df.columns):

            table.cell(
                0,
                col_idx
            ).text = str(column_name)

        # DATA
        for row_idx, row in enumerate(
            self.df.values,
            start=1
        ):

            for col_idx, value in enumerate(row):

                table.cell(
                    row_idx,
                    col_idx
                ).text = str(value)

        # =========================
        # EXPORT
        # =========================

        os.makedirs(
            "exports",
            exist_ok=True
        )

        filename = (
            f"exports/report_"
            f"{uuid.uuid4()}.pptx"
        )

        prs.save(filename)

        return filename
